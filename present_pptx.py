from pptx import Presentation
from pptx.util import Pt,Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import random

def create_presentation(data):
    
  prs = Presentation()
  prs.slide_width = Inches(13.33)
  prs.slide_height = Inches(7.5)
  title_slide_layout  = prs.slide_layouts[1]

  img_path_list = [
     "images\\rm222batch3-kul-15.jpg",
     "images\\simple-blue-white-background-with-text-space.jpg",
     "images\\dvn6.jpg",
     "images\\chris-appano--sTwytNnqWw-unsplash.jpg"
     ]
  img_path = random.choice(img_path_list)


  side_img_path = ["images\\gemini-native-image0.png","images\\gemini-native-image1.png","images\\gemini-native-image2.png"]
  target_slide_indexes = [1,3,5]
  for i in range(len(data)):
      slide = prs.slides.add_slide(title_slide_layout)
      title = slide.shapes.title
      subtitle = slide.placeholders[1]
      
      #backgroundImage
      prs.slide_width,prs.slide_height
      pic = slide.shapes.add_picture(img_path,0,0,width=prs.slide_width,height=prs.slide_height)
      slide.shapes._spTree.remove(pic._element)
      slide.shapes._spTree.insert(2,pic._element)

      #Title
      title.text =  data[i]["title"]
      title.text_frame.paragraphs[0].font.size = Pt(36)
      title.text_frame.paragraphs[0].font.bold = True
      title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


      #Content
      subtitle.text = data[i]["content"]
      p = subtitle.text_frame.paragraphs[0]
      p.font.size = Pt(20)
      if img_path == ".venv\\images\\chris-appano--sTwytNnqWw-unsplash.jpg":
         p.font.color.rgb = RGBColor(250,250,250)
      else:
         p.font.color.rgb = RGBColor(50,50,50)
      p.alignment = PP_ALIGN.JUSTIFY

      subtitle.text_frame.margin_left = Inches(0.2)
      subtitle.text_frame.margin_right = Inches(0.2)


      #Image
      

      img_width = Inches(6)
      img_height = Inches(2)
      img_left = prs.slide_width - img_width - Inches(5)
      img_top = prs.slide_height - img_height - Inches(0.5)

      if i in target_slide_indexes:
            index_in_list = target_slide_indexes.index(i)
            side_img = side_img_path[index_in_list]
            slide.shapes.add_picture(side_img,img_left,img_top,img_width,img_height)

  prs.save("presentation.pptx")

